#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
FES PowerPoint Generator Module
Creates daily generation forecast presentations for DAM meetings
with enhanced charts and trading-relevant analytics
"""

import pandas as pd
from datetime import datetime, timedelta
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import numpy as np


class ForecastPresentationGenerator:
    def __init__(self):
        """Initialize the presentation generator"""
        self.prs = None
        self.trading_date = None
        self.gen_forecast_df = None
        self.d2_forecast_df = None
        self.d3_forecast_df = None  # Saturday
        self.d4_forecast_df = None  # Sunday
        self.d5_forecast_df = None  # Monday
        self.is_friday_presentation = False
        self.chart_dir = Path.cwd() / "output"
        self.chart_dir.mkdir(exist_ok=True)
        
    def load_forecast_data(self, trading_date_str, force_friday_mode=False):
        """
        Load D-1 and D-2 forecast files if they exist
        For Friday presentations, also load D-3, D-4, D-5 (Saturday, Sunday, Monday)
        
        Args:
            trading_date_str: Trading date in DD/MM/YYYY format
            force_friday_mode: If True, force Friday mode regardless of day
        """
        self.trading_date = datetime.strptime(trading_date_str, "%d/%m/%Y")
        
        # Check if this is a Friday presentation (trading day is Friday)
        # Friday = 4 in weekday() where Monday=0
        self.is_friday_presentation = (self.trading_date.weekday() == 4) or force_friday_mode
        
        if self.is_friday_presentation:
            print(f"[PPT] Friday presentation detected - will load weekend forecasts")
        
        # Build file paths for D-1 forecast
        year = str(self.trading_date.year)
        month = self.trading_date.strftime("%B")
        day_str_d1 = self.trading_date.strftime("%d.%m.%Y")
        
        forecast_dir = Path(rf"I:\Daily Generation Forecasts\Daily Generation to Submit\{year}\{month}")
        
        # Load D-1 forecast
        d1_file = forecast_dir / f"Generation Forecast {day_str_d1} D-1.xlsx"
        print(f"[PPT] Loading D-1 forecast: {d1_file}")
        
        if d1_file.exists():
            self.gen_forecast_df = pd.read_excel(d1_file)
            self.gen_forecast_df["DateTime"] = pd.to_datetime(self.gen_forecast_df["DateTime"], format="%d/%m/%Y %H:%M")
        else:
            raise FileNotFoundError(f"D-1 forecast file not found: {d1_file}")
        
        # Try to load D-2 forecast - same trading day but created yesterday
        day_str_d2 = self.trading_date.strftime("%d.%m.%Y")
        d2_file = forecast_dir / f"Generation Forecast {day_str_d2} D-2.xlsx"
        
        if d2_file.exists():
            print(f"[PPT] Loading D-2 forecast: {d2_file}")
            self.d2_forecast_df = pd.read_excel(d2_file)
            self.d2_forecast_df["DateTime"] = pd.to_datetime(self.d2_forecast_df["DateTime"], format="%d/%m/%Y %H:%M")
        else:
            print(f"[PPT] D-2 forecast not found (will skip comparison)")
            self.d2_forecast_df = None
        
        # If Friday presentation, load weekend forecasts (D-3, D-4, D-5)
        if self.is_friday_presentation:
            self._load_weekend_forecasts()
    
    def _load_weekend_forecasts(self):
        """Load D-3 (Saturday), D-4 (Sunday), D-5 (Monday) forecasts for Friday presentations"""
        
        # D-3 = Saturday (trading_date + 1 day)
        saturday = self.trading_date + timedelta(days=1)
        self._load_forecast_for_date(saturday, "D-3", "d3_forecast_df")
        
        # D-4 = Sunday (trading_date + 2 days)
        sunday = self.trading_date + timedelta(days=2)
        self._load_forecast_for_date(sunday, "D-4", "d4_forecast_df")
        
        # D-5 = Monday (trading_date + 3 days)
        monday = self.trading_date + timedelta(days=3)
        self._load_forecast_for_date(monday, "D-5", "d5_forecast_df")
    
    def _load_forecast_for_date(self, date_obj, forecast_label, attr_name):
        """Helper to load a forecast file for a specific date"""
        year = str(date_obj.year)
        month = date_obj.strftime("%B")
        day_str = date_obj.strftime("%d.%m.%Y")
        
        forecast_dir = Path(rf"I:\Daily Generation Forecasts\Daily Generation to Submit\{year}\{month}")
        forecast_file = forecast_dir / f"Generation Forecast {day_str} {forecast_label}.xlsx"
        
        if forecast_file.exists():
            print(f"[PPT] Loading {forecast_label} forecast: {forecast_file}")
            df = pd.read_excel(forecast_file)
            df["DateTime"] = pd.to_datetime(df["DateTime"], format="%d/%m/%Y %H:%M")
            setattr(self, attr_name, df)
        else:
            print(f"[PPT] {forecast_label} forecast not found: {forecast_file}")
            setattr(self, attr_name, None)
    
    def create_line_chart(self):
        """Create line chart showing all generation sources"""
        chart_file = self.chart_dir / "gen_forecast_line.png"
        
        fig, ax = plt.subplots(figsize=(16, 9))
        
        # Plot each generation source
        sources = {
            'Meteo ROI (MW)': '#1f77b4',
            'Meteo NI (MW)': '#ff7f0e',
            'Meteo TB (MW)': '#2ca02c',
            'Meteo CK (MW)': '#d62728',
            'Meteo LD (MW)': '#9467bd',
            'Meteo CD (MW)': '#8c564b',
            'Naïve Nonwind (MW)': '#e377c2',
            'Self-forecast (MW)': '#7f7f7f',
            'Meteo DT (MW)': '#bcbd22',
            'Meteo MUR (MW)': '#17becf'
        }
        
        for col, color in sources.items():
            if col in self.gen_forecast_df.columns:
                ax.plot(self.gen_forecast_df['DateTime'], 
                       self.gen_forecast_df[col], 
                       label=col, 
                       linewidth=2.5,
                       color=color)
        
        ax.set_xlabel('Time', fontsize=14, fontweight='bold')
        ax.set_ylabel('MW', fontsize=14, fontweight='bold')
        ax.set_title(f'D-1 Generation Forecast {self.trading_date.strftime("%d/%m/%Y")}', 
                     fontsize=18, fontweight='bold', pad=20)
        
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%d/%m/%Y %H:%M'))
        ax.xaxis.set_major_locator(mdates.HourLocator(interval=2))
        plt.xticks(rotation=45, ha='right')
        
        ax.grid(True, alpha=0.3, linestyle='--')
        ax.legend(loc='upper left', fontsize=10, ncol=2)
        
        plt.tight_layout()
        plt.savefig(chart_file, dpi=300, bbox_inches='tight')
        plt.close()
        
        return chart_file
    
    def create_stacked_area_chart(self):
        """Create stacked area chart"""
        chart_file = self.chart_dir / "gen_forecast_stacked.png"
        
        fig, ax = plt.subplots(figsize=(16, 9))
        
        # Prepare data for stacking
        stack_cols = []
        labels = []
        colors = []
        
        col_mapping = {
            'Meteo ROI (MW)': ('Meteo ROI (MW)', '#1f77b4'),
            'Meteo NI (MW)': ('Meteo NI (MW)', '#ff7f0e'),
            'Meteo TB (MW)': ('Meteo TB (MW)', '#2ca02c'),
            'Meteo CK (MW)': ('Meteo CK (MW)', '#d62728'),
            'Meteo LD (MW)': ('Meteo LD (MW)', '#9467bd'),
            'Meteo CD (MW)': ('Meteo CD (MW)', '#8c564b'),
            'Naïve Nonwind (MW)': ('Naïve Nonwind (MW)', '#e377c2'),
            'Self-forecast (MW)': ('Self-forecast (MW)', '#7f7f7f'),
            'Meteo DT (MW)': ('Meteo DT (MW)', '#bcbd22'),
            'Meteo MUR (MW)': ('Meteo MUR (MW)', '#17becf')
        }
        
        for col, (label, color) in col_mapping.items():
            if col in self.gen_forecast_df.columns:
                stack_cols.append(self.gen_forecast_df[col].values)
                labels.append(label)
                colors.append(color)
        
        times = self.gen_forecast_df['DateTime']
        
        ax.stackplot(times, *stack_cols, labels=labels, colors=colors, alpha=0.8)
        
        ax.set_xlabel('Time', fontsize=14, fontweight='bold')
        ax.set_ylabel('MW', fontsize=14, fontweight='bold')
        ax.set_title(f'D-1 Generation Forecast {self.trading_date.strftime("%d/%m/%Y")}', 
                     fontsize=18, fontweight='bold', pad=20)
        
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%d/%m/%Y %H:%M'))
        ax.xaxis.set_major_locator(mdates.HourLocator(interval=2))
        plt.xticks(rotation=45, ha='right')
        
        ax.grid(True, alpha=0.3, linestyle='--', axis='y')
        ax.legend(loc='upper left', fontsize=10, ncol=2)
        
        plt.tight_layout()
        plt.savefig(chart_file, dpi=300, bbox_inches='tight')
        plt.close()
        
        return chart_file
    
    def create_d1_vs_d2_chart(self):
        """Create D-1 vs D-2 comparison chart"""
        if self.d2_forecast_df is None:
            return None
        
        chart_file = self.chart_dir / "d1_vs_d2_comparison.png"
        
        fig, ax = plt.subplots(figsize=(16, 9))
        
        # Calculate total generation for each forecast
        gen_cols = ['Meteo ROI (MW)', 'Meteo NI (MW)', 'Meteo TB (MW)', 
                    'Meteo CK (MW)', 'Meteo LD (MW)', 'Meteo CD (MW)',
                    'Naïve Nonwind (MW)', 'Self-forecast (MW)', 'Meteo DT (MW)', 'Meteo MUR (MW)']
        
        # Filter to only existing columns
        d1_cols = [col for col in gen_cols if col in self.gen_forecast_df.columns]
        d2_cols = [col for col in gen_cols if col in self.d2_forecast_df.columns]
        
        d1_total = self.gen_forecast_df[d1_cols].sum(axis=1)
        d2_total = self.d2_forecast_df[d2_cols].sum(axis=1)
        
        # Calculate totals in MWh
        d1_mwh = d1_total.sum() / 2  # Half-hourly to MWh
        d2_mwh = d2_total.sum() / 2
        diff_mwh = d1_mwh - d2_mwh
        
        ax.plot(self.gen_forecast_df['DateTime'], d1_total, 
               label='D-1', linewidth=3, color='#1f77b4', marker='o', markersize=3)
        ax.plot(self.d2_forecast_df['DateTime'], d2_total, 
               label='D-2', linewidth=3, color='#ff7f0e', marker='o', markersize=3)
        
        # Add data labels on points
        for i in range(0, len(d1_total), 4):  # Every 4th point (2 hours)
            ax.text(self.gen_forecast_df['DateTime'].iloc[i], 
                   d1_total.iloc[i], 
                   f'{d1_total.iloc[i]:.1f}', 
                   fontsize=8, ha='center', va='bottom', color='#1f77b4')
        
        ax.set_xlabel('Time', fontsize=14, fontweight='bold')
        ax.set_ylabel('MW', fontsize=14, fontweight='bold')
        
        title_text = f'D-1 change versus D-2\n'
        title_text += f'D-1: {d1_mwh:.1f} MWh\n'
        title_text += f'D-2: {d2_mwh:.1f} MWh\n'
        title_text += f'Diff: {diff_mwh:+.1f} MWh'
        
        ax.set_title(title_text, fontsize=16, fontweight='bold', pad=20)
        
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%d/%m/%Y %H:%M'))
        ax.xaxis.set_major_locator(mdates.HourLocator(interval=2))
        plt.xticks(rotation=45, ha='right')
        
        ax.grid(True, alpha=0.3, linestyle='--')
        ax.legend(loc='upper right', fontsize=12)
        
        plt.tight_layout()
        plt.savefig(chart_file, dpi=300, bbox_inches='tight')
        plt.close()
        
        return chart_file
    
    def create_weekend_forecasts_chart(self):
        """Create chart showing Saturday, Sunday, Monday forecasts for Friday presentations"""
        if not self.is_friday_presentation:
            return None
        
        chart_file = self.chart_dir / "weekend_forecasts.png"
        
        fig, (ax1, ax2, ax3) = plt.subplots(3, 1, figsize=(16, 18))
        
        gen_cols = ['Meteo ROI (MW)', 'Meteo NI (MW)', 'Meteo TB (MW)', 
                    'Meteo CK (MW)', 'Meteo LD (MW)', 'Meteo CD (MW)',
                    'Naïve Nonwind (MW)', 'Self-forecast (MW)', 'Meteo DT (MW)', 'Meteo MUR (MW)']
        
        # Saturday (D-3)
        if self.d3_forecast_df is not None:
            available_cols = [col for col in gen_cols if col in self.d3_forecast_df.columns]
            saturday_total = self.d3_forecast_df[available_cols].sum(axis=1)
            saturday_date = self.d3_forecast_df['DateTime'].iloc[0].strftime("%d/%m/%Y")
            
            ax1.plot(self.d3_forecast_df['DateTime'], saturday_total, 
                    linewidth=3, color='#3498db', label=f'Saturday {saturday_date}')
            ax1.fill_between(self.d3_forecast_df['DateTime'], saturday_total, 0, alpha=0.3, color='#3498db')
            ax1.set_ylabel('MW', fontsize=12, fontweight='bold')
            ax1.set_title(f'Saturday Forecast - {saturday_date}', fontsize=14, fontweight='bold')
            ax1.grid(True, alpha=0.3, linestyle='--')
            ax1.legend(loc='upper right')
            ax1.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M'))
        else:
            ax1.text(0.5, 0.5, 'Saturday forecast not available', 
                    ha='center', va='center', fontsize=14, transform=ax1.transAxes)
        
        # Sunday (D-4)
        if self.d4_forecast_df is not None:
            available_cols = [col for col in gen_cols if col in self.d4_forecast_df.columns]
            sunday_total = self.d4_forecast_df[available_cols].sum(axis=1)
            sunday_date = self.d4_forecast_df['DateTime'].iloc[0].strftime("%d/%m/%Y")
            
            ax2.plot(self.d4_forecast_df['DateTime'], sunday_total, 
                    linewidth=3, color='#e74c3c', label=f'Sunday {sunday_date}')
            ax2.fill_between(self.d4_forecast_df['DateTime'], sunday_total, 0, alpha=0.3, color='#e74c3c')
            ax2.set_ylabel('MW', fontsize=12, fontweight='bold')
            ax2.set_title(f'Sunday Forecast - {sunday_date}', fontsize=14, fontweight='bold')
            ax2.grid(True, alpha=0.3, linestyle='--')
            ax2.legend(loc='upper right')
            ax2.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M'))
        else:
            ax2.text(0.5, 0.5, 'Sunday forecast not available', 
                    ha='center', va='center', fontsize=14, transform=ax2.transAxes)
        
        # Monday (D-5)
        if self.d5_forecast_df is not None:
            available_cols = [col for col in gen_cols if col in self.d5_forecast_df.columns]
            monday_total = self.d5_forecast_df[available_cols].sum(axis=1)
            monday_date = self.d5_forecast_df['DateTime'].iloc[0].strftime("%d/%m/%Y")
            
            ax3.plot(self.d5_forecast_df['DateTime'], monday_total, 
                    linewidth=3, color='#2ecc71', label=f'Monday {monday_date}')
            ax3.fill_between(self.d5_forecast_df['DateTime'], monday_total, 0, alpha=0.3, color='#2ecc71')
            ax3.set_xlabel('Time', fontsize=12, fontweight='bold')
            ax3.set_ylabel('MW', fontsize=12, fontweight='bold')
            ax3.set_title(f'Monday Forecast - {monday_date}', fontsize=14, fontweight='bold')
            ax3.grid(True, alpha=0.3, linestyle='--')
            ax3.legend(loc='upper right')
            ax3.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M'))
        else:
            ax3.text(0.5, 0.5, 'Monday forecast not available', 
                    ha='center', va='center', fontsize=14, transform=ax3.transAxes)
        
        plt.tight_layout()
        plt.savefig(chart_file, dpi=300, bbox_inches='tight')
        plt.close()
        
        return chart_file
    
    def create_forecasting_context_chart(self):
        """Create forecasting context chart showing last 3 days of D-1 forecasts"""
        chart_file = self.chart_dir / "forecasting_context.png"
        
        fig, ax = plt.subplots(figsize=(16, 9))
        
        gen_cols = ['Meteo ROI (MW)', 'Meteo NI (MW)', 'Meteo TB (MW)', 
                    'Meteo CK (MW)', 'Meteo LD (MW)', 'Meteo CD (MW)',
                    'Naïve Nonwind (MW)', 'Self-forecast (MW)', 'Meteo DT (MW)', 'Meteo MUR (MW)']
        
        # Try to load D-1 forecasts for the past 3 days
        all_forecasts = []
        
        for days_back in range(3):
            try:
                past_date = self.trading_date - timedelta(days=days_back)
                year = str(past_date.year)
                month = past_date.strftime("%B")
                day_str = past_date.strftime("%d.%m.%Y")
                
                forecast_dir = Path(rf"I:\Daily Generation Forecasts\Daily Generation to Submit\{year}\{month}")
                forecast_file = forecast_dir / f"Generation Forecast {day_str} D-1.xlsx"
                
                if forecast_file.exists():
                    df = pd.read_excel(forecast_file)
                    df["DateTime"] = pd.to_datetime(df["DateTime"], format="%d/%m/%Y %H:%M")
                    
                    available_cols = [col for col in gen_cols if col in df.columns]
                    df['Total'] = df[available_cols].sum(axis=1)
                    
                    all_forecasts.append(df[['DateTime', 'Total']])
                    print(f"[PPT] Loaded forecast for {day_str}")
            except Exception as e:
                print(f"[PPT] Could not load forecast for {days_back} days back: {e}")
        
        if all_forecasts:
            # Combine all forecasts
            combined_df = pd.concat(all_forecasts, ignore_index=True)
            combined_df = combined_df.sort_values('DateTime')
            
            ax.plot(combined_df['DateTime'], combined_df['Total'], 
                   label='Complete DA Generation Forecasts', 
                   linewidth=3, color='#e74c3c')
        else:
            # Fallback to just current forecast
            available_cols = [col for col in gen_cols if col in self.gen_forecast_df.columns]
            total_gen = self.gen_forecast_df[available_cols].sum(axis=1)
            ax.plot(self.gen_forecast_df['DateTime'], total_gen, 
                   label='Complete DA Generation Forecasts', 
                   linewidth=3, color='#e74c3c')
        
        ax.set_xlabel('Time', fontsize=14, fontweight='bold')
        ax.set_ylabel('MW', fontsize=14, fontweight='bold')
        ax.set_title('Forecasting Context', fontsize=18, fontweight='bold', pad=20)
        
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%d/%m/%Y %H:%M'))
        ax.xaxis.set_major_locator(mdates.DayLocator())
        plt.xticks(rotation=45, ha='right')
        
        ax.grid(True, alpha=0.3, linestyle='--')
        ax.legend(loc='upper right', fontsize=12)
        
        plt.tight_layout()
        plt.savefig(chart_file, dpi=300, bbox_inches='tight')
        plt.close()
        
        return chart_file
    
    def create_presentation(self, trading_date_str, gu_chart_path=None, su_chart_path=None, send_email=False, force_friday_mode=False):
        """
        Create complete PowerPoint presentation
        For Friday presentations, includes Saturday, Sunday, Monday forecasts
        
        Args:
            trading_date_str: Trading date in DD/MM/YYYY format
            gu_chart_path: Path to GU bid chart (optional)
            su_chart_path: Path to SU bid chart (optional)
            send_email: If True, send email with D-1 forecast attached
            force_friday_mode: If True, force Friday mode regardless of day
        
        Returns:
            Path to saved presentation
        """
        print("\n" + "="*70)
        print("CREATING POWERPOINT PRESENTATION")
        print("="*70)
        
        # Load forecast data
        self.load_forecast_data(trading_date_str, force_friday_mode)
        
        # Create new presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Slide 1: Title
        self._create_title_slide()
        
        # Slide 2: Line chart
        print("[PPT] Creating line chart...")
        line_chart = self.create_line_chart()
        self._add_chart_slide("D-1 Generation Forecast", line_chart)
        
        # Slide 3: Stacked area chart
        print("[PPT] Creating stacked area chart...")
        stacked_chart = self.create_stacked_area_chart()
        self._add_chart_slide("D-1 Generation Forecast", stacked_chart)
        
        # Slide 4: D-1 vs D-2 comparison (if D-2 exists)
        if self.d2_forecast_df is not None:
            print("[PPT] Creating D-1 vs D-2 comparison...")
            comparison_chart = self.create_d1_vs_d2_chart()
            if comparison_chart:
                self._add_chart_slide("D-1 change versus D-2", comparison_chart)
        
        # Slide 5: Forecasting context
        print("[PPT] Creating forecasting context...")
        context_chart = self.create_forecasting_context_chart()
        self._add_chart_slide("Forecasting Context", context_chart)
        
        # Friday Special: Weekend forecasts (Saturday, Sunday, Monday)
        if self.is_friday_presentation:
            print("[PPT] Creating weekend forecasts (D-3, D-4, D-5)...")
            weekend_chart = self.create_weekend_forecasts_chart()
            if weekend_chart:
                self._add_chart_slide("Weekend Forecasts (Sat-Sun-Mon)", weekend_chart)
        
        # Slide 6: Outages (placeholder for now)
        self._create_outages_slide()
        
        # Slide 7: GU Bid Chart (if provided)
        if gu_chart_path and Path(gu_chart_path).exists():
            print("[PPT] Adding Murley GU bid chart...")
            self._add_chart_slide("Murley GU504260 Bid Submission", gu_chart_path)
        
        # Slide 8: SU Bid Chart (if provided)
        if su_chart_path and Path(su_chart_path).exists():
            print("[PPT] Adding Supply Unit bid chart...")
            self._add_chart_slide("Supply Unit SU400130 Bid Submission", su_chart_path)
        
        # Save presentation
        ppt_path = self._save_presentation(trading_date_str)
        
        print("="*70)
        print(f"[PPT] Presentation created: {ppt_path}")
        if self.is_friday_presentation:
            print(f"[PPT] Friday presentation - included weekend forecasts")
        print("="*70 + "\n")
        
        # Send email if requested
        if send_email:
            self.send_forecast_email(trading_date_str)
        
        return ppt_path
    
    def send_forecast_email(self, trading_date_str):
        """
        Send email to trading team with D-1 forecast attached
        
        Args:
            trading_date_str: Trading date in DD/MM/YYYY format
        """
        print("\n" + "="*70)
        print("SENDING EMAIL TO TRADING TEAM")
        print("="*70)
        
        try:
            import win32com.client
        except ImportError:
            print("[ERROR] pywin32 not installed. Install with: pip install pywin32")
            return
        
        try:
            # Get forecast file path
            date_obj = datetime.strptime(trading_date_str, "%d/%m/%Y")
            year = str(date_obj.year)
            month = date_obj.strftime("%B")
            day_str = date_obj.strftime("%d.%m.%Y")
            
            forecast_dir = Path(rf"I:\Daily Generation Forecasts\Daily Generation to Submit\{year}\{month}")
            forecast_file = forecast_dir / f"Generation Forecast {day_str} D-1.xlsx"
            
            if not forecast_file.exists():
                print(f"[ERROR] Forecast file not found: {forecast_file}")
                return
            
            # Calculate totals
            gen_cols = ['Meteo ROI (MW)', 'Meteo NI (MW)', 'Meteo TB (MW)', 
                        'Meteo CK (MW)', 'Meteo LD (MW)', 'Meteo CD (MW)',
                        'Naïve Nonwind (MW)', 'Self-forecast (MW)', 'Meteo DT (MW)', 'Meteo MUR (MW)']
            
            available_cols = [col for col in gen_cols if col in self.gen_forecast_df.columns]
            d1_total_mwh = (self.gen_forecast_df[available_cols].sum(axis=1).sum()) / 2
            
            # Build email body
            day_name = date_obj.strftime("%A")
            formatted_date = date_obj.strftime("%d/%m/%Y")
            
            email_body = f"Hi Trading,<br><br>"
            email_body += f"Please see attached D-1 forecast for {formatted_date}<br><br>"
            email_body += f"Updated Forecast D-1: {d1_total_mwh:.1f} MWh<br>"
            
            if self.d2_forecast_df is not None:
                d2_cols = [col for col in gen_cols if col in self.d2_forecast_df.columns]
                d2_total_mwh = (self.d2_forecast_df[d2_cols].sum(axis=1).sum()) / 2
                diff_mwh = d1_total_mwh - d2_total_mwh
                email_body += f"Previous Forecast D-2: {d2_total_mwh:.1f} MWh<br>"
                email_body += f"<br>Diff: {diff_mwh:+.1f} MWh<br>"
            
            email_body += f"<br>No dial applied.<br>"
            email_body += f"<br>Kind Regards,<br>Renewables Team"
            
            # Create Outlook email
            outlook = win32com.client.Dispatch("Outlook.Application")
            mail = outlook.CreateItem(0)
            
            # Find the correct account
            for account in outlook.Session.Accounts:
                if account.DisplayName == "manasvin.mahajan@flogas.ie":
                    mail._oleobj_.Invoke(*(64209, 0, 8, 0, account))
                    break
            
            # Set email properties
            mail.To = "isemtrading@flogas.ie"
            mail.CC = "generator_forecasting@flogas.ie"
            mail.Subject = f"ISEM D-1 Generation Volumes {day_name} {date_obj.strftime('%d/%m/%Y')}"
            mail.HTMLBody = email_body
            mail.Attachments.Add(str(forecast_file))
            
            # Send
            mail.Send()
            
            print("[EMAIL] ✓ Email sent successfully")
            print(f"[EMAIL] To: isemtrading@flogas.ie")
            print(f"[EMAIL] CC: generator_forecasting@flogas.ie")
            print(f"[EMAIL] Subject: ISEM D-1 Generation Volumes {day_name} {date_obj.strftime('%d/%m/%Y')}")
            print(f"[EMAIL] Attachment: {forecast_file.name}")
            print("="*70 + "\n")
            
        except Exception as e:
            print(f"[ERROR] Email sending failed: {e}")
            print("="*70 + "\n")
    
    def _create_title_slide(self):
        """Create title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"D-1 Generation Forecast {self.trading_date.strftime('%d/%m/%Y')}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        title_para.alignment = PP_ALIGN.CENTER
    
    def _add_chart_slide(self, title, chart_path):
        """Add a slide with a chart"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add chart image
        slide.shapes.add_picture(
            str(chart_path),
            Inches(0.5), Inches(1.2),
            width=Inches(9), height=Inches(5.8)
        )
    
    def _create_outages_slide(self):
        """Create outages slide (placeholder)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Outages"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add placeholder text
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        content_frame = content_box.text_frame
        content_frame.text = "No outages for specified trading day(s)"
        content_para = content_frame.paragraphs[0]
        content_para.font.size = Pt(24)
        content_para.alignment = PP_ALIGN.CENTER
    
    def _save_presentation(self, trading_date_str):
        """Save presentation to file"""
        date_obj = datetime.strptime(trading_date_str, "%d/%m/%Y")
        year = str(date_obj.year)
        month_num = str(date_obj.month)
        month_name = date_obj.strftime("%B")
        day_str = date_obj.strftime("%d.%m.%Y")
        
        # Create output directory
        ppt_dir = Path(rf"V:\Renewables\Presentations and Learning\4) Daily Forecast meetings\{year}\{month_num}) {month_name}")
        ppt_dir.mkdir(parents=True, exist_ok=True)
        
        # Save file
        ppt_file = ppt_dir / f"Generation_Forecast_{day_str.replace('.', '_')}.pptx"
        self.prs.save(str(ppt_file))
        
        return ppt_file


# Standalone function for easy calling
def generate_forecast_presentation(trading_date_str, gu_chart_path=None, su_chart_path=None, send_email=False, force_friday_mode=False):
    """
    Generate forecast presentation and optionally send email
    For Friday presentations, automatically includes weekend forecasts (Sat, Sun, Mon)
    
    Args:
        trading_date_str: Trading date in DD/MM/YYYY format
        gu_chart_path: Path to GU bid chart PNG (optional)
        su_chart_path: Path to SU bid chart PNG (optional)
        send_email: If True, send email to trading team with D-1 forecast attached
        force_friday_mode: If True, force Friday mode (load weekend forecasts) regardless of day
    
    Returns:
        Path to saved presentation
    """
    generator = ForecastPresentationGenerator()
    return generator.create_presentation(trading_date_str, gu_chart_path, su_chart_path, send_email, force_friday_mode)


if __name__ == "__main__":
    # Test
    ppt_path = generate_forecast_presentation("23/01/2026")
    print(f"Presentation saved: {ppt_path}")